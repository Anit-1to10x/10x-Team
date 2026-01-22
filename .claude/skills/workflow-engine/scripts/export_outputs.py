#!/usr/bin/env python3
"""
Output Exporter - Generates PDF, PPT, and other output formats
Developed by Team 10x.in
"""

import json
import os
import sys
from datetime import datetime
from pathlib import Path

# Try to import optional dependencies
try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image, PageBreak
    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from PIL import Image, ImageDraw, ImageFont
    HAS_PILLOW = True
except ImportError:
    HAS_PILLOW = False

ROOT_DIR = Path(__file__).resolve().parent.parent.parent.parent.parent


class OutputExporter:
    """Exports workflow results to various formats."""

    def __init__(self, workflow_path: str):
        self.workflow_path = Path(workflow_path)
        self.workflow_dir = self.workflow_path.parent
        self.workflow = None

    def load_workflow(self):
        """Load workflow from file."""
        with open(self.workflow_path, "r") as f:
            self.workflow = json.load(f)
        return self.workflow

    def export_json(self) -> str:
        """Export workflow as JSON."""
        output_file = self.workflow_dir / "final_output.json"

        # Compile all step outputs
        compiled = {
            "workflow_id": self.workflow["workflow_id"],
            "name": self.workflow["name"],
            "description": self.workflow["description"],
            "executed_at": self.workflow["execution"].get("completed_at"),
            "status": self.workflow["status"],
            "user_inputs": self.workflow["user_inputs"].get("answers", {}),
            "steps": [],
            "developed_by": "Team 10x.in"
        }

        for step in self.workflow["steps"]:
            step_data = {
                "step_id": step["step_id"],
                "name": step["name"],
                "skill": step["skill"],
                "status": step.get("status", "unknown"),
                "outputs": step.get("outputs", {})
            }
            compiled["steps"].append(step_data)

        with open(output_file, "w") as f:
            json.dump(compiled, f, indent=2)

        print(f"JSON exported: {output_file}")
        return str(output_file)

    def export_pdf(self) -> str:
        """Export workflow as PDF."""
        if not HAS_REPORTLAB:
            print("reportlab not installed. Install with: pip install reportlab")
            return None

        output_file = self.workflow_dir / "final_output.pdf"

        doc = SimpleDocTemplate(
            str(output_file),
            pagesize=A4,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=72
        )

        styles = getSampleStyleSheet()

        # Custom styles
        title_style = ParagraphStyle(
            'Title',
            parent=styles['Heading1'],
            fontSize=24,
            spaceAfter=30,
            textColor=colors.HexColor('#1a1a2e')
        )

        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=16,
            spaceAfter=12,
            textColor=colors.HexColor('#16213e')
        )

        body_style = ParagraphStyle(
            'CustomBody',
            parent=styles['Normal'],
            fontSize=11,
            spaceAfter=8,
            textColor=colors.HexColor('#333333')
        )

        # Build PDF content
        story = []

        # Title
        story.append(Paragraph(self.workflow["name"], title_style))
        story.append(Spacer(1, 12))

        # Metadata
        story.append(Paragraph("Workflow Details", heading_style))
        story.append(Paragraph(f"<b>Workflow ID:</b> {self.workflow['workflow_id']}", body_style))
        story.append(Paragraph(f"<b>Status:</b> {self.workflow['status']}", body_style))
        story.append(Paragraph(f"<b>Created:</b> {self.workflow['created_at']}", body_style))

        if self.workflow["execution"].get("completed_at"):
            story.append(Paragraph(f"<b>Completed:</b> {self.workflow['execution']['completed_at']}", body_style))

        story.append(Spacer(1, 20))

        # Description
        story.append(Paragraph("Description", heading_style))
        story.append(Paragraph(self.workflow["description"], body_style))
        story.append(Spacer(1, 20))

        # User Inputs
        if self.workflow["user_inputs"].get("answers"):
            story.append(Paragraph("User Inputs", heading_style))
            for key, value in self.workflow["user_inputs"]["answers"].items():
                story.append(Paragraph(f"<b>{key}:</b> {value}", body_style))
            story.append(Spacer(1, 20))

        # Steps
        story.append(Paragraph("Workflow Steps", heading_style))

        for step in self.workflow["steps"]:
            story.append(Paragraph(
                f"<b>Step {step['step_id']}:</b> {step['name']}",
                body_style
            ))
            story.append(Paragraph(f"Skill: {step['skill']}", body_style))
            story.append(Paragraph(f"Status: {step.get('status', 'pending')}", body_style))

            if step.get("outputs"):
                story.append(Paragraph(f"Output: {json.dumps(step['outputs'], indent=2)[:500]}...", body_style))

            story.append(Spacer(1, 12))

        story.append(Spacer(1, 30))

        # Footer
        story.append(Paragraph("─" * 50, body_style))
        story.append(Paragraph("Developed by Team 10x.in | https://10x.in", body_style))

        # Build PDF
        doc.build(story)
        print(f"PDF exported: {output_file}")
        return str(output_file)

    def export_ppt(self) -> str:
        """Export workflow as PowerPoint presentation."""
        if not HAS_PPTX:
            print("python-pptx not installed. Install with: pip install python-pptx")
            return None

        output_file = self.workflow_dir / "final_output.pptx"

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        title_slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(title_slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = self.workflow["name"]
        title_para.font.size = Pt(44)
        title_para.font.bold = True

        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = self.workflow["description"]
        subtitle_para.font.size = Pt(24)

        # Add footer
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
        footer_frame = footer_box.text_frame
        footer_para = footer_frame.paragraphs[0]
        footer_para.text = "Developed by Team 10x.in"
        footer_para.font.size = Pt(14)

        # Steps slides
        for step in self.workflow["steps"]:
            slide = prs.slides.add_slide(title_slide_layout)

            # Step title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = f"Step {step['step_id']}: {step['name']}"
            title_para.font.size = Pt(32)
            title_para.font.bold = True

            # Step details
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
            content_frame = content_box.text_frame

            # Skill
            p = content_frame.paragraphs[0]
            p.text = f"Skill: {step['skill']}"
            p.font.size = Pt(18)

            # Status
            p = content_frame.add_paragraph()
            p.text = f"Status: {step.get('status', 'pending')}"
            p.font.size = Pt(18)

            # Outputs
            if step.get("outputs"):
                p = content_frame.add_paragraph()
                p.text = ""
                p.font.size = Pt(12)

                p = content_frame.add_paragraph()
                p.text = "Outputs:"
                p.font.size = Pt(18)
                p.font.bold = True

                output_text = json.dumps(step["outputs"], indent=2)
                for line in output_text.split("\n")[:10]:  # Limit to 10 lines
                    p = content_frame.add_paragraph()
                    p.text = line
                    p.font.size = Pt(12)

        # Summary slide
        slide = prs.slides.add_slide(title_slide_layout)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Summary"
        title_para.font.size = Pt(32)
        title_para.font.bold = True

        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
        content_frame = content_box.text_frame

        p = content_frame.paragraphs[0]
        p.text = f"Workflow ID: {self.workflow['workflow_id']}"
        p.font.size = Pt(18)

        p = content_frame.add_paragraph()
        p.text = f"Total Steps: {len(self.workflow['steps'])}"
        p.font.size = Pt(18)

        p = content_frame.add_paragraph()
        p.text = f"Status: {self.workflow['status']}"
        p.font.size = Pt(18)

        if self.workflow["execution"].get("completed_at"):
            p = content_frame.add_paragraph()
            p.text = f"Completed: {self.workflow['execution']['completed_at']}"
            p.font.size = Pt(18)

        # Save
        prs.save(output_file)
        print(f"PPT exported: {output_file}")
        return str(output_file)

    def export_png(self) -> str:
        """Export workflow as PNG diagram."""
        if not HAS_PILLOW:
            print("Pillow not installed. Install with: pip install Pillow")
            return None

        output_file = self.workflow_dir / "workflow_diagram.png"

        # Calculate dimensions
        step_count = len(self.workflow["steps"])
        width = 1200
        header_height = 200
        step_height = 120
        footer_height = 80
        height = header_height + (step_count * step_height) + footer_height

        # Create image
        img = Image.new('RGB', (width, height), color='#FFFFFF')
        draw = ImageDraw.Draw(img)

        # Try to load a font, fallback to default
        try:
            title_font = ImageFont.truetype("arial.ttf", 32)
            heading_font = ImageFont.truetype("arial.ttf", 20)
            body_font = ImageFont.truetype("arial.ttf", 14)
        except:
            title_font = ImageFont.load_default()
            heading_font = ImageFont.load_default()
            body_font = ImageFont.load_default()

        # Colors
        primary_color = '#1a1a2e'
        secondary_color = '#16213e'
        accent_color = '#0f3460'
        success_color = '#00a878'
        pending_color = '#f4a261'
        border_color = '#e0e0e0'

        # Draw header background
        draw.rectangle([0, 0, width, header_height], fill='#f8f9fa')

        # Draw title
        draw.text((50, 40), self.workflow["name"], fill=primary_color, font=title_font)

        # Draw workflow ID and status
        draw.text((50, 90), f"Workflow ID: {self.workflow['workflow_id']}", fill=secondary_color, font=body_font)
        status = self.workflow["status"]
        status_color = success_color if status == "completed" else pending_color
        draw.text((50, 115), f"Status: {status}", fill=status_color, font=body_font)
        draw.text((50, 140), f"Steps: {step_count}", fill=secondary_color, font=body_font)

        # Draw steps
        y_offset = header_height + 20
        for i, step in enumerate(self.workflow["steps"]):
            # Step box
            box_x = 50
            box_y = y_offset
            box_width = width - 100
            box_height = step_height - 20

            # Determine step color based on status
            step_status = step.get("status", "pending")
            if step_status == "completed":
                box_color = '#e8f5e9'
                border = success_color
            elif step_status == "in_progress":
                box_color = '#fff3e0'
                border = pending_color
            else:
                box_color = '#f5f5f5'
                border = border_color

            # Draw step box
            draw.rectangle([box_x, box_y, box_x + box_width, box_y + box_height], fill=box_color, outline=border, width=2)

            # Step number circle
            circle_x = box_x + 30
            circle_y = box_y + box_height // 2
            circle_radius = 18
            draw.ellipse([circle_x - circle_radius, circle_y - circle_radius, circle_x + circle_radius, circle_y + circle_radius], fill=accent_color)
            draw.text((circle_x - 6, circle_y - 8), str(i + 1), fill='white', font=body_font)

            # Step name
            draw.text((circle_x + 40, box_y + 15), step["name"], fill=primary_color, font=heading_font)

            # Step skill and status
            draw.text((circle_x + 40, box_y + 45), f"Skill: {step['skill']}", fill=secondary_color, font=body_font)
            draw.text((circle_x + 40, box_y + 65), f"Status: {step_status}", fill=status_color if step_status != "pending" else secondary_color, font=body_font)

            # Draw connector line
            if i < step_count - 1:
                line_x = width // 2
                line_y_start = box_y + box_height
                line_y_end = box_y + step_height
                draw.line([(line_x, line_y_start), (line_x, line_y_end)], fill=border_color, width=2)
                # Arrow
                draw.polygon([(line_x - 6, line_y_end - 8), (line_x + 6, line_y_end - 8), (line_x, line_y_end)], fill=border_color)

            y_offset += step_height

        # Draw footer
        footer_y = height - footer_height + 20
        draw.rectangle([0, height - footer_height, width, height], fill='#f8f9fa')
        draw.text((50, footer_y), "Developed by Team 10x.in | https://10x.in", fill=secondary_color, font=body_font)

        # Save
        img.save(output_file)
        print(f"PNG exported: {output_file}")
        return str(output_file)

    def export_svg(self) -> str:
        """Export workflow as SVG diagram."""
        output_file = self.workflow_dir / "workflow_diagram.svg"

        # Calculate dimensions
        step_count = len(self.workflow["steps"])
        width = 1200
        header_height = 200
        step_height = 120
        footer_height = 80
        height = header_height + (step_count * step_height) + footer_height

        # SVG content
        svg_content = f'''<?xml version="1.0" encoding="UTF-8"?>
<svg width="{width}" height="{height}" xmlns="http://www.w3.org/2000/svg">
  <defs>
    <style>
      .title {{ font-family: Arial, sans-serif; font-size: 32px; font-weight: bold; fill: #1a1a2e; }}
      .heading {{ font-family: Arial, sans-serif; font-size: 20px; font-weight: bold; fill: #1a1a2e; }}
      .body {{ font-family: Arial, sans-serif; font-size: 14px; fill: #16213e; }}
      .small {{ font-family: Arial, sans-serif; font-size: 12px; fill: white; }}
      .success {{ fill: #00a878; }}
      .pending {{ fill: #f4a261; }}
    </style>
  </defs>

  <!-- Header Background -->
  <rect x="0" y="0" width="{width}" height="{header_height}" fill="#f8f9fa"/>

  <!-- Title -->
  <text x="50" y="60" class="title">{self._escape_xml(self.workflow["name"])}</text>

  <!-- Metadata -->
  <text x="50" y="100" class="body">Workflow ID: {self.workflow["workflow_id"]}</text>
  <text x="50" y="125" class="body {'success' if self.workflow["status"] == 'completed' else 'pending'}">Status: {self.workflow["status"]}</text>
  <text x="50" y="150" class="body">Steps: {step_count}</text>

'''

        y_offset = header_height + 20
        for i, step in enumerate(self.workflow["steps"]):
            step_status = step.get("status", "pending")
            box_x = 50
            box_y = y_offset
            box_width = width - 100
            box_height = step_height - 20

            # Colors based on status
            if step_status == "completed":
                box_fill = "#e8f5e9"
                border_color = "#00a878"
            elif step_status == "in_progress":
                box_fill = "#fff3e0"
                border_color = "#f4a261"
            else:
                box_fill = "#f5f5f5"
                border_color = "#e0e0e0"

            # Step box
            svg_content += f'''
  <!-- Step {i + 1} -->
  <rect x="{box_x}" y="{box_y}" width="{box_width}" height="{box_height}" fill="{box_fill}" stroke="{border_color}" stroke-width="2" rx="8"/>

  <!-- Step Number -->
  <circle cx="{box_x + 30}" cy="{box_y + box_height // 2}" r="18" fill="#0f3460"/>
  <text x="{box_x + 30}" y="{box_y + box_height // 2 + 5}" class="small" text-anchor="middle">{i + 1}</text>

  <!-- Step Content -->
  <text x="{box_x + 70}" y="{box_y + 30}" class="heading">{self._escape_xml(step["name"])}</text>
  <text x="{box_x + 70}" y="{box_y + 55}" class="body">Skill: {step["skill"]}</text>
  <text x="{box_x + 70}" y="{box_y + 75}" class="body">Status: {step_status}</text>
'''

            # Connector line
            if i < step_count - 1:
                line_x = width // 2
                line_y_start = box_y + box_height
                line_y_end = box_y + step_height - 5
                svg_content += f'''
  <line x1="{line_x}" y1="{line_y_start}" x2="{line_x}" y2="{line_y_end}" stroke="#e0e0e0" stroke-width="2"/>
  <polygon points="{line_x - 6},{line_y_end - 8} {line_x + 6},{line_y_end - 8} {line_x},{line_y_end}" fill="#e0e0e0"/>
'''

            y_offset += step_height

        # Footer
        footer_y = height - footer_height
        svg_content += f'''
  <!-- Footer -->
  <rect x="0" y="{footer_y}" width="{width}" height="{footer_height}" fill="#f8f9fa"/>
  <text x="50" y="{footer_y + 35}" class="body">Developed by Team 10x.in | https://10x.in</text>

</svg>'''

        with open(output_file, "w", encoding="utf-8") as f:
            f.write(svg_content)

        print(f"SVG exported: {output_file}")
        return str(output_file)

    def _escape_xml(self, text: str) -> str:
        """Escape XML special characters."""
        if not text:
            return ""
        return (str(text)
                .replace("&", "&amp;")
                .replace("<", "&lt;")
                .replace(">", "&gt;")
                .replace('"', "&quot;")
                .replace("'", "&apos;"))

    def export_markdown(self) -> str:
        """Export workflow as Markdown."""
        output_file = self.workflow_dir / "final_output.md"

        content = f"""# {self.workflow['name']}

{self.workflow['description']}

---

## Workflow Details

- **Workflow ID:** {self.workflow['workflow_id']}
- **Status:** {self.workflow['status']}
- **Created:** {self.workflow['created_at']}
- **Completed:** {self.workflow['execution'].get('completed_at', 'N/A')}

---

## User Inputs

"""
        if self.workflow["user_inputs"].get("answers"):
            for key, value in self.workflow["user_inputs"]["answers"].items():
                content += f"- **{key}:** {value}\n"

        content += """
---

## Workflow Steps

"""
        for step in self.workflow["steps"]:
            content += f"""### Step {step['step_id']}: {step['name']}

- **Skill:** {step['skill']}
- **Action:** {step.get('action', 'execute')}
- **Status:** {step.get('status', 'pending')}

"""
            if step.get("outputs"):
                content += "**Outputs:**\n```json\n"
                content += json.dumps(step["outputs"], indent=2)
                content += "\n```\n\n"

        content += """---

*Developed by Team 10x.in | https://10x.in*
"""

        with open(output_file, "w") as f:
            f.write(content)

        print(f"Markdown exported: {output_file}")
        return str(output_file)

    def export_all(self, formats: list = None) -> dict:
        """Export to all requested formats."""
        if formats is None:
            formats = ["json", "md", "png"]  # Default formats

        results = {}

        for fmt in formats:
            if fmt == "json":
                results["json"] = self.export_json()
            elif fmt == "pdf":
                results["pdf"] = self.export_pdf()
            elif fmt == "ppt" or fmt == "pptx":
                results["ppt"] = self.export_ppt()
            elif fmt == "md" or fmt == "markdown":
                results["md"] = self.export_markdown()
            elif fmt == "png":
                results["png"] = self.export_png()
            elif fmt == "svg":
                results["svg"] = self.export_svg()

        return results


def main():
    if len(sys.argv) < 2:
        print("Usage: python export_outputs.py <workflow_path> [formats]")
        print("Formats: json, pdf, ppt, md, png, svg (comma-separated)")
        print("Example: python export_outputs.py workflow.json json,pdf,ppt,png,svg")
        sys.exit(1)

    workflow_path = sys.argv[1]
    formats = sys.argv[2].split(",") if len(sys.argv) > 2 else ["json", "md"]

    if not os.path.exists(workflow_path):
        print(f"Workflow not found: {workflow_path}")
        sys.exit(1)

    exporter = OutputExporter(workflow_path)
    exporter.load_workflow()
    results = exporter.export_all(formats)

    print("\n--- EXPORT RESULTS ---")
    print(json.dumps(results, indent=2))


if __name__ == "__main__":
    main()
